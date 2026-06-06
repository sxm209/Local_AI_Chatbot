from __future__ import annotations

import csv
import hashlib
import html
import logging
from dataclasses import dataclass
from html.parser import HTMLParser
from pathlib import Path
from typing import Any

LOGGER = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".txt",
    ".md",
    ".markdown",
    ".csv",
    ".xlsx",
    ".pptx",
    ".html",
    ".htm",
}


@dataclass(frozen=True)
class ExtractedSection:
    text: str
    metadata: dict[str, Any]


@dataclass(frozen=True)
class ExtractedDocument:
    path: Path
    title: str
    file_type: str
    sha256: str
    modified_at: float
    sections: list[ExtractedSection]
    warning: str | None = None


class _HTMLTextParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []

    def handle_data(self, data: str) -> None:
        value = data.strip()
        if value:
            self.parts.append(value)

    def text(self) -> str:
        return "\n".join(self.parts)


def is_supported(path: Path) -> bool:
    return path.suffix.lower() in SUPPORTED_EXTENSIONS


def file_sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for block in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest()


def discover_files(paths: list[str], recursive: bool = True) -> tuple[list[Path], list[dict[str, str]]]:
    files: list[Path] = []
    warnings: list[dict[str, str]] = []
    for raw in paths:
        candidate = Path(raw).expanduser()
        if candidate.is_file():
            if is_supported(candidate):
                files.append(candidate)
            else:
                warnings.append({"path": str(candidate), "message": "Unsupported file type"})
        elif candidate.is_dir():
            iterator = candidate.rglob("*") if recursive else candidate.glob("*")
            for child in iterator:
                if child.is_file():
                    if is_supported(child):
                        files.append(child)
                    else:
                        warnings.append({"path": str(child), "message": "Unsupported file type"})
        else:
            warnings.append({"path": str(candidate), "message": "Path does not exist"})
    return files, warnings


def extract_document(path: Path) -> ExtractedDocument:
    path = path.expanduser().resolve()
    suffix = path.suffix.lower()
    stat = path.stat()
    sha = file_sha256(path)
    warning = None

    try:
        if suffix in {".txt", ".md", ".markdown"}:
            sections = _extract_plain(path)
        elif suffix in {".html", ".htm"}:
            sections = _extract_html(path)
        elif suffix == ".csv":
            sections = _extract_csv(path)
        elif suffix == ".pdf":
            sections = _extract_pdf(path)
        elif suffix == ".docx":
            sections = _extract_docx(path)
        elif suffix == ".xlsx":
            sections = _extract_xlsx(path)
        elif suffix == ".pptx":
            sections = _extract_pptx(path)
        else:
            raise ValueError(f"Unsupported file type: {suffix}")
    except ImportError as exc:
        warning = f"Missing parser dependency: {exc.name or str(exc)}"
        sections = []
    except Exception as exc:
        LOGGER.info("extract_failed path=%s error=%s", path, type(exc).__name__)
        warning = f"Could not extract text: {exc}"
        sections = []

    return ExtractedDocument(
        path=path,
        title=path.name,
        file_type=suffix.lstrip("."),
        sha256=sha,
        modified_at=stat.st_mtime,
        sections=sections,
        warning=warning,
    )


def _extract_plain(path: Path) -> list[ExtractedSection]:
    text = path.read_text(encoding="utf-8", errors="replace")
    return [ExtractedSection(text=text, metadata={"section": "body"})]


def _extract_html(path: Path) -> list[ExtractedSection]:
    raw = path.read_text(encoding="utf-8", errors="replace")
    parser = _HTMLTextParser()
    parser.feed(raw)
    text = html.unescape(parser.text())
    return [ExtractedSection(text=text, metadata={"section": "html"})]


def _extract_csv(path: Path) -> list[ExtractedSection]:
    rows: list[str] = []
    with path.open("r", encoding="utf-8", errors="replace", newline="") as handle:
        reader = csv.reader(handle)
        for index, row in enumerate(reader, start=1):
            rows.append(f"Row {index}: " + " | ".join(cell.strip() for cell in row))
    return [ExtractedSection(text="\n".join(rows), metadata={"section": "csv"})]


def _extract_pdf(path: Path) -> list[ExtractedSection]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    sections: list[ExtractedSection] = []
    for index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        sections.append(ExtractedSection(text=text, metadata={"page": index}))
    return sections


def _extract_docx(path: Path) -> list[ExtractedSection]:
    import docx

    document = docx.Document(str(path))
    text = "\n".join(paragraph.text for paragraph in document.paragraphs if paragraph.text.strip())
    return [ExtractedSection(text=text, metadata={"section": "document"})]


def _extract_xlsx(path: Path) -> list[ExtractedSection]:
    import openpyxl

    workbook = openpyxl.load_workbook(path, read_only=True, data_only=True)
    sections: list[ExtractedSection] = []
    for sheet in workbook.worksheets:
        rows: list[str] = []
        for index, row in enumerate(sheet.iter_rows(values_only=True), start=1):
            values = ["" if cell is None else str(cell) for cell in row]
            if any(value.strip() for value in values):
                rows.append(f"Row {index}: " + " | ".join(values))
        sections.append(ExtractedSection(text="\n".join(rows), metadata={"sheet": sheet.title}))
    return sections


def _extract_pptx(path: Path) -> list[ExtractedSection]:
    from pptx import Presentation

    presentation = Presentation(str(path))
    sections: list[ExtractedSection] = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text)
        sections.append(ExtractedSection(text="\n".join(texts), metadata={"slide": index}))
    return sections
